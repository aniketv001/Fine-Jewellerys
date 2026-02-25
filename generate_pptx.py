#!/usr/bin/env python3
"""
Fine Jewellery's — PowerPoint Presentation Generator
Generates a 15-slide .pptx file with screenshots and content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Paths ──────────────────────────────────────────────────────────────────
ARTIFACTS = "/Users/aniket/.gemini/antigravity/brain/e1480d75-4e98-4529-875c-b7236ac34c1e"
OUTPUT = "/Users/aniket/Desktop/Fine Jewellery Website/Fine_Jewellery_Presentation.pptx"

SCREENSHOTS = {
    "homepage":       os.path.join(ARTIFACTS, "slide_homepage_1771445921278.png"),
    "products":       os.path.join(ARTIFACTS, "slide_products_1771445932814.png"),
    "product_detail": os.path.join(ARTIFACTS, "slide_product_detail_retry_1771445960252.png"),
    "cart":           os.path.join(ARTIFACTS, "slide_cart_1771445971627.png"),
    "wishlist":       os.path.join(ARTIFACTS, "slide_wishlist_1771445981728.png"),
    "auth":           os.path.join(ARTIFACTS, "slide_auth_1771445991221.png"),
    "checkout":       os.path.join(ARTIFACTS, "slide_checkout_1771446000797.png"),
}

# ── Colors ─────────────────────────────────────────────────────────────────
GOLD        = RGBColor(0xC9, 0xA8, 0x4C)
GOLD_LIGHT  = RGBColor(0xE8, 0xC9, 0x7A)
DARK        = RGBColor(0x0D, 0x0D, 0x0D)
DARK2       = RGBColor(0x1A, 0x1A, 0x1A)
DARK3       = RGBColor(0x25, 0x25, 0x25)
WHITE       = RGBColor(0xFA, 0xFA, 0xFA)
MUTED       = RGBColor(0x88, 0x88, 0x88)
RED         = RGBColor(0xEE, 0x55, 0x55)

# ── Slide size: 16:9 widescreen ────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=Pt(14), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font_name="Calibri", italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_gold_line(slide, left, top, width=Inches(0.8)):
    rect = add_rect(slide, left, top, width, Pt(3), fill_color=GOLD)
    return rect

def slide_number_label(slide, num, label):
    text = f"{num:02d}  —  {label.upper()}"
    add_text(slide, text, Inches(0.6), Inches(0.35), Inches(6), Inches(0.3),
             font_size=Pt(9), color=GOLD, bold=True, font_name="Calibri")

def section_title(slide, title, top=Inches(0.7)):
    add_text(slide, title, Inches(0.6), top, Inches(12), Inches(0.9),
             font_size=Pt(36), bold=True, color=WHITE, font_name="Georgia")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, DARK)

# Gold glow rectangle (decorative)
add_rect(slide, Inches(2), Inches(1.2), Inches(9.33), Inches(5.1), fill_color=RGBColor(0x1A, 0x12, 0x00))

# Diamond emoji / icon area
add_text(slide, "💎", Inches(5.5), Inches(0.8), Inches(2.33), Inches(0.8),
         font_size=Pt(40), align=PP_ALIGN.CENTER, color=WHITE)

# Main title
add_text(slide, "Fine Jewellery's",
         Inches(1), Inches(1.6), Inches(11.33), Inches(1.5),
         font_size=Pt(60), bold=True, color=GOLD_LIGHT,
         align=PP_ALIGN.CENTER, font_name="Georgia")

# Subtitle
add_text(slide, "E-COMMERCE PLATFORM  —  PROJECT PRESENTATION",
         Inches(1), Inches(3.1), Inches(11.33), Inches(0.5),
         font_size=Pt(14), color=MUTED, align=PP_ALIGN.CENTER, bold=False)

# Gold line
add_gold_line(slide, Inches(6.1), Inches(3.7), Inches(1.1))

# Stats row
stats = [("15", "SLIDES"), ("7+", "PAGES"), ("Full", "STACK"), ("Live", "DEMO")]
for i, (num, lbl) in enumerate(stats):
    x = Inches(2.5 + i * 2.1)
    add_text(slide, num, x, Inches(4.1), Inches(1.8), Inches(0.6),
             font_size=Pt(28), bold=True, color=GOLD, align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text(slide, lbl, x, Inches(4.65), Inches(1.8), Inches(0.3),
             font_size=Pt(9), color=MUTED, align=PP_ALIGN.CENTER)

# Presenter line
add_text(slide, "Presented by  Aniket   ·   Fine Jewellery's E-Commerce Project   ·   2026",
         Inches(1), Inches(5.5), Inches(11.33), Inches(0.4),
         font_size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Project Overview
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, DARK2)
slide_number_label(slide, 2, "Overview")
section_title(slide, "Project Overview")
add_gold_line(slide, Inches(0.6), Inches(1.65))

# Description
desc = ("Fine Jewellery's is a full-stack luxury e-commerce platform built to showcase and sell "
        "premium handcrafted jewellery. The platform offers a seamless shopping experience with "
        "a modern, elegant UI inspired by high-end jewellery brands.")
add_text(slide, desc, Inches(0.6), Inches(1.9), Inches(6.2), Inches(1.2),
         font_size=Pt(12), color=RGBColor(0xBB, 0xBB, 0xBB))

# Bullet points
bullets = [
    "✦  Complete e-commerce flow from browsing to checkout",
    "✦  User authentication with JWT tokens",
    "✦  Persistent wishlist synced to MongoDB",
    "✦  Shopping cart with quantity management",
    "✦  Responsive design for all screen sizes",
    "✦  BIS Hallmarked jewellery catalogue",
]
for i, b in enumerate(bullets):
    add_text(slide, b, Inches(0.6), Inches(3.15 + i * 0.42), Inches(6.2), Inches(0.4),
             font_size=Pt(11), color=RGBColor(0xCC, 0xCC, 0xCC))

# Right column — info cards
info = [
    ("PROJECT TYPE", "Full-Stack Web Application"),
    ("DOMAIN", "Luxury E-Commerce / Jewellery"),
    ("TARGET USERS", "Jewellery shoppers across India"),
    ("DEPLOYMENT", "Local Dev · Cloud-ready"),
]
for i, (label, val) in enumerate(info):
    y = Inches(1.9 + i * 1.3)
    add_rect(slide, Inches(7.2), y, Inches(5.5), Inches(1.1),
             fill_color=DARK3, line_color=RGBColor(0x40, 0x35, 0x15))
    add_text(slide, label, Inches(7.4), y + Pt(8), Inches(5), Inches(0.3),
             font_size=Pt(9), color=GOLD, bold=True)
    add_text(slide, val, Inches(7.4), y + Inches(0.4), Inches(5), Inches(0.5),
             font_size=Pt(14), color=WHITE, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Tech Stack
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, DARK)
slide_number_label(slide, 3, "Technology")
section_title(slide, "Tech Stack")
add_gold_line(slide, Inches(0.6), Inches(1.65))

tech = [
    ("⚛️ React 18 + TypeScript", "Component-based UI with full type safety and modern hooks"),
    ("⚡ Vite", "Lightning-fast build tool and dev server with HMR"),
    ("🎨 Tailwind CSS", "Utility-first CSS framework for rapid, consistent styling"),
    ("🟢 Node.js + Express", "RESTful API backend with middleware-based architecture"),
    ("🍃 MongoDB + Mongoose", "NoSQL database with schema validation and Atlas cloud hosting"),
    ("🔐 JWT Authentication", "Stateless auth with bcrypt password hashing"),
    ("🔀 React Router v6", "Client-side routing with protected routes and navigation"),
    ("🧩 Radix UI + shadcn", "Accessible component primitives with custom theming"),
    ("🔔 Sonner Toasts", "Beautiful animated toast notifications for user feedback"),
]

cols = 3
rows = 3
for i, (name, desc) in enumerate(tech):
    col = i % cols
    row = i // cols
    x = Inches(0.5 + col * 4.27)
    y = Inches(2.0 + row * 1.6)
    add_rect(slide, x, y, Inches(4.0), Inches(1.4),
             fill_color=DARK3, line_color=RGBColor(0x40, 0x35, 0x15))
    add_text(slide, name, x + Inches(0.15), y + Inches(0.1), Inches(3.7), Inches(0.45),
             font_size=Pt(12), bold=True, color=GOLD_LIGHT)
    add_text(slide, desc, x + Inches(0.15), y + Inches(0.55), Inches(3.7), Inches(0.7),
             font_size=Pt(10), color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDES 4–10 — Screenshot slides
# ══════════════════════════════════════════════════════════════════════════════
screenshot_slides = [
    (4, "Homepage",         "homepage",       [
        "Auto-rotating hero carousel with 3 banners",
        "Top announcement bar with offers",
        "Smart search bar with placeholder hints",
        "Category navigation: Gold, Silver, Men, Women",
        "Wishlist & Cart count badges in header",
        "Delivery pincode selector",
    ]),
    (5, "Product Listing",  "products",       [
        "Grid layout with hover animations",
        "Filter by category, material, price range",
        "NEW & BESTSELLER badges on cards",
        "Quick 'Add to Cart' on hover",
        "Heart icon for instant wishlist toggle",
        "Star ratings and review counts",
    ]),
    (6, "Product Detail",   "product_detail", [
        "High-res product image gallery (4 views)",
        "Breadcrumb navigation",
        "Quantity selector with +/- controls",
        "Add to Cart & Wishlist buttons",
        "Buy Now → direct checkout",
        "Related products section",
    ]),
    (7, "Shopping Cart",    "cart",           [
        "Persistent cart saved in localStorage",
        "Quantity update with +/- controls",
        "Remove individual items",
        "Real-time price calculation",
        "Order summary with subtotal",
        "Proceed to Checkout CTA",
    ]),
    (8, "Wishlist Feature", "wishlist",       [
        "Synced to MongoDB when logged in",
        "localStorage fallback when not logged in",
        "Optimistic UI — instant heart toggle",
        "Move to Cart with one click",
        "Wishlist count badge in header",
        "Persists across sessions",
    ]),
    (9, "Login & Sign Up",  "auth",           [
        "JWT-based stateless authentication",
        "bcrypt password hashing (12 rounds)",
        "Token stored in localStorage",
        "Auto-login on page refresh",
        "Form validation with error messages",
        "Protected routes for authenticated users",
    ]),
    (10, "Checkout Page",   "checkout",       [
        "Shipping address form",
        "Order summary with item list",
        "Price breakdown (subtotal, shipping, total)",
        "Payment method selection",
        "Free shipping on all orders",
        "Order confirmation flow",
    ]),
]

for slide_num, title, key, bullets in screenshot_slides:
    slide = blank_slide(prs)
    fill_bg(slide, DARK)

    # Left panel background
    add_rect(slide, 0, 0, Inches(5.5), H, fill_color=DARK2)

    # Slide number label
    slide_number_label(slide, slide_num, title)

    # Title
    add_text(slide, title, Inches(0.5), Inches(0.75), Inches(4.8), Inches(1.0),
             font_size=Pt(32), bold=True, color=WHITE, font_name="Georgia")

    # Gold line
    add_gold_line(slide, Inches(0.5), Inches(1.8))

    # Bullets
    for i, b in enumerate(bullets):
        add_text(slide, f"◆  {b}", Inches(0.5), Inches(2.1 + i * 0.72), Inches(4.8), Inches(0.6),
                 font_size=Pt(11), color=RGBColor(0xCC, 0xCC, 0xCC))

    # Screenshot on right
    img_path = SCREENSHOTS.get(key)
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path,
                                 Inches(5.6), Inches(0.1),
                                 Inches(7.6), Inches(7.3))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Backend Architecture
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, DARK2)
slide_number_label(slide, 11, "Architecture")
section_title(slide, "Backend Architecture")
add_gold_line(slide, Inches(0.6), Inches(1.65))

# Frontend box
add_rect(slide, Inches(0.4), Inches(2.0), Inches(4.0), Inches(4.0),
         fill_color=DARK3, line_color=RGBColor(0x40, 0x35, 0x15))
add_text(slide, "🖥️  Frontend (React)", Inches(0.55), Inches(2.1), Inches(3.7), Inches(0.5),
         font_size=Pt(14), bold=True, color=GOLD_LIGHT)
fe_items = ["React 18 + TypeScript + Vite", "AuthContext — JWT token management",
            "CartContext — cart & wishlist state", "React Router v6 — client routing",
            "Tailwind CSS + Radix UI", "Sonner toast notifications",
            "localStorage for cart persistence"]
for i, item in enumerate(fe_items):
    add_text(slide, f"▸  {item}", Inches(0.6), Inches(2.65 + i * 0.47), Inches(3.6), Inches(0.4),
             font_size=Pt(10), color=RGBColor(0xBB, 0xBB, 0xBB))

# Arrow
add_text(slide, "REST API\n──────────\nHTTP / JSON\nBearer Token",
         Inches(4.55), Inches(3.3), Inches(1.5), Inches(1.2),
         font_size=Pt(10), color=GOLD, align=PP_ALIGN.CENTER)

# Backend box
add_rect(slide, Inches(6.1), Inches(2.0), Inches(4.0), Inches(4.0),
         fill_color=DARK3, line_color=RGBColor(0x40, 0x35, 0x15))
add_text(slide, "⚙️  Backend (Node.js)", Inches(6.25), Inches(2.1), Inches(3.7), Inches(0.5),
         font_size=Pt(14), bold=True, color=GOLD_LIGHT)
be_items = ["Express.js REST API server", "JWT middleware — route protection",
            "bcryptjs — password hashing", "Mongoose ODM — MongoDB models",
            "CORS configured for port 8080", "Routes: /auth, /wishlist, /orders",
            "MongoDB Atlas — cloud database"]
for i, item in enumerate(be_items):
    add_text(slide, f"▸  {item}", Inches(6.3), Inches(2.65 + i * 0.47), Inches(3.6), Inches(0.4),
             font_size=Pt(10), color=RGBColor(0xBB, 0xBB, 0xBB))

# Stats bar
add_rect(slide, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.9),
         fill_color=DARK3, line_color=RGBColor(0x40, 0x35, 0x15))
stats = [("Frontend Port", ":8080"), ("Backend Port", ":5000"),
         ("Database", "MongoDB Atlas"), ("Auth", "JWT + bcrypt")]
for i, (lbl, val) in enumerate(stats):
    x = Inches(1.2 + i * 3.1)
    add_text(slide, lbl, x, Inches(6.25), Inches(2.8), Inches(0.3),
             font_size=Pt(9), color=GOLD, bold=True)
    add_text(slide, val, x, Inches(6.6), Inches(2.8), Inches(0.4),
             font_size=Pt(13), color=WHITE, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Database Schema
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, DARK2)
slide_number_label(slide, 12, "Database")
section_title(slide, "MongoDB Schema Design")
add_gold_line(slide, Inches(0.6), Inches(1.65))

schemas = [
    ("👤 User Model", [
        ("_id", "ObjectId"), ("name", "String *"), ("email", "String * unique"),
        ("phone", "String *"), ("password", "String * hashed"),
        ("role", "user | admin"), ("isActive", "Boolean"), ("createdAt", "Date"),
    ]),
    ("❤️ Wishlist Model", [
        ("_id", "ObjectId"), ("user", "→ User ref"), ("items[]", "Array"),
        ("  productId", "String *"), ("  name", "String *"),
        ("  image", "String *"), ("  price", "Number *"), ("  material", "String"),
    ]),
    ("📦 Order Model", [
        ("_id", "ObjectId"), ("user", "→ User ref"), ("items[]", "Array"),
        ("  productId", "String"), ("  quantity", "Number"),
        ("totalAmount", "Number *"), ("status", "pending|confirmed"),
        ("shippingAddress", "Object"),
    ]),
]

for col, (title, fields) in enumerate(schemas):
    x = Inches(0.4 + col * 4.3)
    # Header
    add_rect(slide, x, Inches(2.0), Inches(4.0), Inches(0.5), fill_color=GOLD)
    add_text(slide, title, x + Inches(0.1), Inches(2.05), Inches(3.8), Inches(0.4),
             font_size=Pt(12), bold=True, color=DARK)
    # Body
    add_rect(slide, x, Inches(2.5), Inches(4.0), Inches(4.5), fill_color=DARK3,
             line_color=RGBColor(0x40, 0x35, 0x15))
    for i, (fname, ftype) in enumerate(fields):
        y = Inches(2.6 + i * 0.47)
        add_text(slide, fname, x + Inches(0.15), y, Inches(2.2), Inches(0.4),
                 font_size=Pt(10), color=WHITE)
        add_text(slide, ftype, x + Inches(2.4), y, Inches(1.5), Inches(0.4),
                 font_size=Pt(9), color=GOLD, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Key Features
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, DARK)
slide_number_label(slide, 13, "Features")
section_title(slide, "Key Features Summary")
add_gold_line(slide, Inches(0.6), Inches(1.65))

features = [
    ("🔐", "Secure Auth", "JWT tokens, bcrypt hashing, protected routes, auto-login"),
    ("🛒", "Smart Cart", "Persistent localStorage cart with quantity management"),
    ("❤️", "Wishlist Sync", "MongoDB-backed wishlist with optimistic UI updates"),
    ("🔍", "Smart Search", "Search across products by name, category, material"),
    ("📱", "Responsive", "Fully responsive design for mobile, tablet, desktop"),
    ("🏷️", "Filters", "Filter by category, material, price range, new arrivals"),
    ("🌙", "Dark Mode", "Elegant dark theme with gold accents throughout"),
    ("🔔", "Notifications", "Animated toast notifications for all user actions"),
]

cols = 4
for i, (icon, name, desc) in enumerate(features):
    col = i % cols
    row = i // cols
    x = Inches(0.4 + col * 3.2)
    y = Inches(2.0 + row * 2.4)
    add_rect(slide, x, y, Inches(3.0), Inches(2.1),
             fill_color=DARK3, line_color=RGBColor(0x40, 0x35, 0x15))
    add_text(slide, icon, x + Inches(0.1), y + Inches(0.1), Inches(2.8), Inches(0.55),
             font_size=Pt(24), align=PP_ALIGN.CENTER)
    add_text(slide, name, x + Inches(0.1), y + Inches(0.7), Inches(2.8), Inches(0.4),
             font_size=Pt(12), bold=True, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x + Inches(0.1), y + Inches(1.1), Inches(2.8), Inches(0.85),
             font_size=Pt(9), color=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Challenges & Solutions
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, DARK2)
slide_number_label(slide, 14, "Challenges")
section_title(slide, "Challenges & Solutions")
add_gold_line(slide, Inches(0.6), Inches(1.65))

challenges = [
    ("MongoDB Atlas DNS SRV records couldn't be resolved by local ISP DNS, causing backend connection failure",
     "Replaced mongodb+srv:// URI with direct shard hostnames discovered via Google DNS (8.8.8.8)"),
    ("Wishlist adding silently failed — async API calls were fired without await, so UI state never updated",
     "Added optimistic UI updates, proper async/await in all handlers, boolean returns, and state revert on failure"),
    ("Rollup native module @rollup/rollup-darwin-arm64 missing due to npm optional dependency bug on Apple Silicon",
     "Deleted node_modules and package-lock.json, then ran a clean npm install to resolve platform-specific binaries"),
]

for i, (challenge, solution) in enumerate(challenges):
    y = Inches(2.0 + i * 1.65)
    # Challenge box
    add_rect(slide, Inches(0.4), y, Inches(5.5), Inches(1.4),
             fill_color=DARK3, line_color=RED)
    add_text(slide, "CHALLENGE", Inches(0.55), y + Inches(0.08), Inches(5), Inches(0.3),
             font_size=Pt(8), bold=True, color=RED)
    add_text(slide, challenge, Inches(0.55), y + Inches(0.4), Inches(5.2), Inches(0.9),
             font_size=Pt(10), color=RGBColor(0xCC, 0xCC, 0xCC))
    # Arrow
    add_text(slide, "→", Inches(6.05), y + Inches(0.5), Inches(0.6), Inches(0.5),
             font_size=Pt(22), color=GOLD, align=PP_ALIGN.CENTER)
    # Solution box
    add_rect(slide, Inches(6.8), y, Inches(6.1), Inches(1.4),
             fill_color=DARK3, line_color=GOLD)
    add_text(slide, "SOLUTION", Inches(6.95), y + Inches(0.08), Inches(5.8), Inches(0.3),
             font_size=Pt(8), bold=True, color=GOLD)
    add_text(slide, solution, Inches(6.95), y + Inches(0.4), Inches(5.8), Inches(0.9),
             font_size=Pt(10), color=RGBColor(0xCC, 0xCC, 0xCC))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, DARK)

# Glow
add_rect(slide, Inches(2), Inches(1), Inches(9.33), Inches(5.5),
         fill_color=RGBColor(0x1A, 0x12, 0x00))

add_text(slide, "💎", Inches(5.5), Inches(0.8), Inches(2.33), Inches(0.9),
         font_size=Pt(44), align=PP_ALIGN.CENTER, color=WHITE)

add_text(slide, "Thank You",
         Inches(1), Inches(1.7), Inches(11.33), Inches(1.4),
         font_size=Pt(64), bold=True, color=GOLD_LIGHT,
         align=PP_ALIGN.CENTER, font_name="Georgia")

add_text(slide, "Questions & Discussion",
         Inches(1), Inches(3.1), Inches(11.33), Inches(0.5),
         font_size=Pt(16), color=MUTED, align=PP_ALIGN.CENTER)

add_gold_line(slide, Inches(6.1), Inches(3.7), Inches(1.1))

add_text(slide, "Fine Jewellery's  —  A Full-Stack Luxury E-Commerce Platform",
         Inches(1), Inches(4.0), Inches(11.33), Inches(0.4),
         font_size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)

# Contact cards
contacts = [("🌐", "LIVE URL", "localhost:8080"), ("⚙️", "API SERVER", "localhost:5000"),
            ("🍃", "DATABASE", "MongoDB Atlas"), ("👤", "DEVELOPER", "Aniket")]
for i, (icon, lbl, val) in enumerate(contacts):
    x = Inches(1.2 + i * 2.8)
    add_rect(slide, x, Inches(4.6), Inches(2.5), Inches(1.5),
             fill_color=DARK3, line_color=RGBColor(0x40, 0x35, 0x15))
    add_text(slide, icon, x, Inches(4.65), Inches(2.5), Inches(0.5),
             font_size=Pt(20), align=PP_ALIGN.CENTER)
    add_text(slide, lbl, x, Inches(5.15), Inches(2.5), Inches(0.3),
             font_size=Pt(8), color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, val, x, Inches(5.45), Inches(2.5), Inches(0.4),
             font_size=Pt(12), bold=True, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"✅ Presentation saved to:\n   {OUTPUT}")
print(f"   Slides: {len(prs.slides)}")
